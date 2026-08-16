#!/usr/bin/env python3
"""Extract plain text from common document formats for the dsh read_document tool."""
import argparse
import html
import re
import sys
import zipfile
from pathlib import Path

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8")

EXT_KINDS = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".xlsx": "xlsx",
    ".xlsm": "xlsx",
    ".pptx": "pptx",
    ".ppsx": "pptx",
    ".odt": "odt",
    ".ods": "ods",
    ".odp": "odp",
    ".epub": "epub",
    ".rtf": "rtf",
    ".txt": "text",
    ".csv": "text",
    ".tsv": "text",
    ".md": "text",
    ".markdown": "text",
    ".json": "text",
    ".jsonl": "text",
    ".log": "text",
    ".xml": "text",
    ".html": "text",
    ".htm": "text",
    ".yaml": "text",
    ".yml": "text",
    ".ini": "text",
    ".toml": "text",
    ".conf": "text",
    ".properties": "text",
}

FORMAT_NAMES = {
    "pdf": "PDF",
    "docx": "Word (.docx)",
    "xlsx": "Excel (.xlsx/.xlsm)",
    "pptx": "PowerPoint (.pptx/.ppsx)",
    "odt": "OpenDocument Text",
    "ods": "OpenDocument Spreadsheet",
    "odp": "OpenDocument Presentation",
    "epub": "EPUB",
    "rtf": "Rich Text Format",
    "text": "plain text",
}


def fail(message):
    print(f"read_document: {message}", file=sys.stderr)
    sys.exit(2)


def clean(text):
    if text is None:
        return ""
    text = str(text).replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+\n", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def strip_tags(xml, paragraph_ends, cell_end="</table:table-cell>"):
    for token in paragraph_ends:
        xml = xml.replace(token, "\n")
    xml = xml.replace(cell_end, "\t")
    text = re.sub(r"<[^>]+>", "", xml)
    return html.unescape(text)


def extract_text_file(path):
    raw = Path(path).read_bytes()
    for encoding in ("utf-8-sig", "utf-16", "gb18030", "latin-1"):
        try:
            return clean(raw.decode(encoding))
        except (UnicodeDecodeError, ValueError):
            continue
    return clean(raw.decode("utf-8", errors="replace"))


def extract_pdf(path):
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    if reader.is_encrypted:
        try:
            reader.decrypt("")
        except Exception:
            pass
    parts = []
    for page in reader.pages:
        try:
            parts.append(page.extract_text() or "")
        except Exception:
            continue
    text = clean("\n\n".join(parts))
    if not text:
        try:
            import pdfplumber

            parts = []
            with pdfplumber.open(str(path)) as pdf:
                for page in pdf.pages:
                    parts.append(page.extract_text() or "")
            text = clean("\n\n".join(parts))
        except Exception:
            pass
    if not text:
        try:
            from pdfminer.high_level import extract_text as pdfminer_text

            text = clean(pdfminer_text(str(path)))
        except Exception:
            pass
    if not text:
        fail("PDF contains no extractable text (it may be a scanned/image PDF requiring OCR)")
    return text


def extract_docx(path):
    from docx import Document
    from docx.oxml.ns import qn

    document = Document(str(path))
    parts = []
    for child in document.element.body.iterchildren():
        if child.tag == qn("w:p"):
            texts = child.findall(".//" + qn("w:t"))
            parts.append("".join(t.text or "" for t in texts))
        elif child.tag == qn("w:tbl"):
            for row in child.findall(qn("w:tr")):
                cells = []
                for cell in row.findall(qn("w:tc")):
                    cell_text = "".join(
                        t.text or "" for t in cell.findall(".//" + qn("w:t"))
                    )
                    cells.append(cell_text)
                parts.append("\t".join(cells))
    return clean("\n".join(parts))


def extract_xlsx(path):
    from openpyxl import load_workbook

    workbook = load_workbook(str(path), read_only=True, data_only=True)
    parts = []
    for sheet in workbook.worksheets:
        parts.append(f"--- Sheet: {sheet.title} ---")
        for row in sheet.iter_rows(values_only=True):
            values = ["" if value is None else str(value) for value in row]
            if any(value.strip() for value in values):
                parts.append("\t".join(values))
    return clean("\n".join(parts))


def extract_pptx(path):
    from pptx import Presentation

    presentation = Presentation(str(path))
    parts = []
    for index, slide in enumerate(presentation.slides, 1):
        parts.append(f"--- Slide {index} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in paragraph.runs)
                    if line.strip():
                        parts.append(line)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    parts.append("\t".join(cell.text for cell in row.cells))
    return clean("\n".join(parts))


def extract_odf(path, kind):
    with zipfile.ZipFile(path) as archive:
        xml = archive.read("content.xml").decode("utf-8", errors="replace")
    if kind == "odt":
        text = strip_tags(
            xml,
            ["</text:p>", "</text:h>", "</text:list-item>", "</table:table-row>"],
        )
    elif kind == "ods":
        text = strip_tags(xml, ["</table:table-row>", "</text:p>"])
    else:
        text = strip_tags(xml, ["</text:p>", "</text:h>", "</text:list-item>"])
    return clean(text)


def extract_epub(path):
    with zipfile.ZipFile(path) as archive:
        names = [
            name
            for name in archive.namelist()
            if name.lower().endswith((".xhtml", ".html", ".htm"))
        ]
        parts = []
        for name in sorted(names):
            raw = archive.read(name).decode("utf-8", errors="replace")
            raw = re.sub(
                r"<(?:script|style)[^>]*>.*?</(?:script|style)>",
                " ",
                raw,
                flags=re.S | re.I,
            )
            raw = re.sub(r"<br[^>]*>", "\n", raw, flags=re.I)
            raw = re.sub(r"</(?:p|div|h[1-6]|li|tr)>", "\n", raw, flags=re.I)
            raw = re.sub(r"<[^>]+>", "", raw)
            part = clean(html.unescape(raw))
            if part:
                parts.append(part)
    return clean("\n\n".join(parts))


def extract_rtf(path):
    raw = Path(path).read_bytes().decode("latin-1", errors="replace")
    raw = re.sub(r"\\u(-?\d+)\??", lambda match: chr(int(match.group(1)) & 0xFFFF), raw)
    raw = re.sub(
        r"\\'([0-9a-fA-F]{2})",
        lambda match: chr(int(match.group(1), 16)),
        raw,
    )
    raw = re.sub(r"\\[a-zA-Z]+-?\d* ?", "", raw)
    raw = re.sub(r"[{}]", "", raw)
    return clean(raw)


def detect_kind(path):
    extension = Path(path).suffix.lower()
    if extension in EXT_KINDS:
        return EXT_KINDS[extension]
    with open(path, "rb") as handle:
        head = handle.read(8)
    if head.startswith(b"%PDF"):
        return "pdf"
    if head.startswith(b"PK\x03\x04"):
        try:
            with zipfile.ZipFile(path) as archive:
                names = set(archive.namelist())
            if any(name.startswith("word/") for name in names):
                return "docx"
            if any(name.startswith("xl/") for name in names):
                return "xlsx"
            if any(name.startswith("ppt/") for name in names):
                return "pptx"
            if "mimetype" in names:
                mimetype = archive.read("mimetype").decode("ascii", errors="ignore")
                if "opendocument.text" in mimetype:
                    return "odt"
                if "opendocument.spreadsheet" in mimetype:
                    return "ods"
                if "opendocument.presentation" in mimetype:
                    return "odp"
                if "epub" in mimetype:
                    return "epub"
        except zipfile.BadZipFile:
            pass
        return None
    return None


HANDLERS = {
    "pdf": extract_pdf,
    "docx": extract_docx,
    "xlsx": extract_xlsx,
    "pptx": extract_pptx,
    "odt": lambda path: extract_odf(path, "odt"),
    "ods": lambda path: extract_odf(path, "ods"),
    "odp": lambda path: extract_odf(path, "odp"),
    "epub": extract_epub,
    "rtf": extract_rtf,
    "text": extract_text_file,
}


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--path", required=True)
    parser.add_argument("--kind", default="auto")
    parser.add_argument("--max-chars", type=int, default=2000000)
    args = parser.parse_args()

    path = Path(args.path)
    if not path.is_file():
        fail(f"file not found: {path}")

    kind = args.kind if args.kind != "auto" else detect_kind(path)
    if kind is None:
        extension = path.suffix.lower()
        if extension in (".doc", ".xls", ".ppt", ".pps"):
            fail(
                f"legacy {extension} binary format is not supported by the built-in reader; "
                "please save the file as .docx/.xlsx/.pptx and try again"
            )
        fail(
            "unsupported file type; supported formats: "
            + ", ".join(sorted(FORMAT_NAMES))
        )

    handler = HANDLERS.get(kind)
    if handler is None:
        fail(f"unsupported file type: {kind}")

    text = handler(path)
    if len(text) > args.max_chars:
        text = text[: args.max_chars]
    sys.stdout.write(text)


if __name__ == "__main__":
    main()
